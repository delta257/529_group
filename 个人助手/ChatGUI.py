import os
import threading
import queue
import time
from datetime import datetime
from dotenv import load_dotenv
import speech_recognition as sr
import pyttsx3
import tkinter as tk
from tkinter import ttk, scrolledtext, messagebox, filedialog
from tkinterdnd2 import DND_FILES, TkinterDnD
from langchain_openai import ChatOpenAI
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationChain
import fitz  # PyMuPDF for PDF
from docx import Document  # python-docx for Word
from pptx import Presentation  # python-pptx for PPT
import base64
from PIL import Image
from io import BytesIO
import cv2
from PIL import ImageTk

class ChatGUI(TkinterDnD.Tk):
    def __init__(self, bot):
        # 初始化队列
        self.queue = queue.Queue()
        self.bot = bot
        self.pending_files = []  # 存储待处理的文件信息
        
        # 调用父类构造函数
        super().__init__()
        
        # 设置主题色
        self.style = ttk.Style()
        
        # 配置基础样式
        self.style.configure('Modern.TFrame', background='#f0f0f0')
        self.style.configure('Modern.TLabel', 
                           background='#f0f0f0', 
                           font=('Microsoft YaHei', 10))
        self.style.configure('Modern.TButton', 
                           font=('Microsoft YaHei', 10),
                           padding=5)
        self.style.configure('Modern.TEntry',
                           font=('Microsoft YaHei', 10),
                           padding=5)
        
        # 配置LabelFrame样式
        self.style.configure('Modern.TLabelframe', 
                           background='#f0f0f0',
                           font=('Microsoft YaHei', 10))
        self.style.configure('Modern.TLabelframe.Label', 
                           background='#f0f0f0',
                           font=('Microsoft YaHei', 10))
        
        self.title("阿kie")
        self._setup_ui()
        self.after(100, self.check_queue)
        
        # 启动持续语音监听
        self.bot.start_listening(self._handle_voice_input)
        
        # 设置文件拖放
        self.drop_target_register(DND_FILES)
        self.dnd_bind('<<Drop>>', self._handle_drop)

    def check_queue(self):
        """检查并处理消息队列"""
        try:
            while True:
                func = self.queue.get_nowait()
                func()
        except queue.Empty:
            pass
        finally:
            self.after(100, self.check_queue)
            
    def _handle_voice_input(self, text):
        """处理语音输入"""
        print(f"收到语音输入: {text}")
        
        # 如果正在处理或说话，跳过这次输入
        if self.bot.is_processing or self.bot.is_speaking:
            print("系统正在处理或说话，跳过此次输入")
            return
            
        # 显示语音输入
        self._thread_safe_display("您", text)
        
        # 处理语音输入
        threading.Thread(
            target=self._process_query,
            args=(text,),
            daemon=True
        ).start()
        
    def _setup_ui(self):
        """构建用户界面"""
        # 设置窗口最小尺寸
        self.minsize(800, 600)
        
        # 创建主框架
        main_frame = ttk.Frame(self, style='Modern.TFrame')
        main_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        # 创建左右分隔框架
        left_frame = ttk.Frame(main_frame, style='Modern.TFrame')
        left_frame.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=(0, 5))
        left_frame.pack_propagate(False)
        left_frame.configure(width=300)  # 设置左侧宽度为总宽度的30%
        
        right_frame = ttk.Frame(main_frame, style='Modern.TFrame')
        right_frame.pack(side=tk.LEFT, fill=tk.BOTH, expand=True, padx=(5, 0))
        
        # 左侧区域：摄像头和文件上传
        self._build_left_section(left_frame)
        
        # 右侧区域：聊天记录和输入控制
        self._build_right_section(right_frame)

    def _build_left_section(self, parent):
        """构建左侧区域（摄像头和文件上传）"""
        # 创建上下分隔框架
        top_frame = ttk.Frame(parent, style='Modern.TFrame')
        top_frame.pack(side=tk.TOP, fill=tk.BOTH, expand=True, pady=(0, 5))
        
        bottom_frame = ttk.Frame(parent, style='Modern.TFrame')
        bottom_frame.pack(side=tk.TOP, fill=tk.BOTH, expand=True, pady=(5, 0))
        
        # 摄像头区域
        camera_frame = ttk.LabelFrame(top_frame, text="摄像头", style='Modern.TLabelframe')
        camera_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
        
        # 摄像头画面显示区域
        self.camera_display = ttk.Label(camera_frame, style='Modern.TLabel')
        self.camera_display.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
        
        # 摄像头控制按钮
        self.camera_button = ttk.Button(
            camera_frame,
            text="开启摄像头",
            command=self._toggle_camera,
            style='Modern.TButton'
        )
        self.camera_button.pack(pady=5)
        
        # 文件上传区域
        upload_frame = ttk.LabelFrame(bottom_frame, text="文件上传", style='Modern.TLabelframe')
        upload_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)
        
        # 创建容器框架用于居中显示
        center_frame = ttk.Frame(upload_frame, style='Modern.TFrame')
        center_frame.pack(expand=True)
        
        # 文件上传按钮
        self.upload_button = ttk.Button(
            center_frame,
            text="选择文件",
            command=self._open_file_dialog,
            style='Modern.TButton'
        )
        self.upload_button.pack(pady=(0, 5))
        
        # 添加拖放提示标签
        self.drop_label = ttk.Label(
            center_frame,
            text="或将文件拖放到此处",
            wraplength=250,
            style='Modern.TLabel'
        )
        self.drop_label.pack()
        
        # 文件处理状态显示区域
        self.file_status_frame = ttk.Frame(upload_frame, style='Modern.TFrame')
        self.file_status_frame.pack(fill=tk.X, padx=5, pady=5)
        
        # 文件处理状态标签
        self.file_status_label = ttk.Label(
            self.file_status_frame,
            text="",
            wraplength=250,
            style='Modern.TLabel'
        )
        self.file_status_label.pack()

    def _build_right_section(self, parent):
        """构建右侧区域（聊天记录和输入控制）"""
        # 聊天记录区
        self._build_chat_area(parent)
        
        # 输入控制区
        self._build_input_controls(parent)

    def _build_chat_area(self, parent):
        """构建聊天记录显示区域"""
        # 创建聊天区域框架
        chat_frame = ttk.Frame(parent, style='Modern.TFrame')
        chat_frame.pack(expand=True, fill=tk.BOTH)
        
        # 聊天记录显示区域
        self.chat_area = scrolledtext.ScrolledText(
            chat_frame,
            wrap=tk.WORD,
            state='disabled',
            font=('Microsoft YaHei', 10),
            spacing3=5,  # 段落间距
            bg='#ffffff',  # 白色背景
            fg='#333333',  # 深灰色文字
            padx=10,
            pady=10
        )
        self.chat_area.pack(expand=True, fill=tk.BOTH, padx=5, pady=5)
        
        # 初始化欢迎信息
        self._thread_safe_display("系统", "请通过语音或文字开始对话")

    def _build_input_controls(self, parent):
        """构建输入控制区域"""
        control_frame = ttk.Frame(parent, style='Modern.TFrame')
        control_frame.pack(fill=tk.X, pady=5)
        
        # 文本输入框
        self.input_text = ttk.Entry(control_frame, style='Modern.TEntry')
        self.input_text.pack(side=tk.LEFT, expand=True, fill=tk.X, padx=(0, 5))
        
        # 发送按钮
        self.send_button = ttk.Button(
            control_frame,
            text="发送",
            command=self._on_text_input,
            style='Modern.TButton'
        )
        self.send_button.pack(side=tk.LEFT, padx=5)
        
        # 语音播报开关按钮
        self.tts_button = ttk.Button(
            control_frame,
            text="语音播报: 开",
            command=self._toggle_tts,
            style='Modern.TButton'
        )
        self.tts_button.pack(side=tk.LEFT, padx=5)
        
        # 绑定回车键
        self.input_text.bind('<Return>', lambda e: self._on_text_input())

    def _thread_safe_display(self, sender, message):
        """线程安全的消息显示方法"""
        self.queue.put(lambda: self._display_message(sender, message))

    def _display_message(self, sender, message):
        """实际执行消息显示"""
        self.chat_area.config(state='normal')
        timestamp = datetime.now().strftime("%H:%M:%S")
        
        # 根据发送者设置不同的颜色
        if sender == "系统":
            color = "#666666"  # 灰色
        elif sender == "您":
            color = "#2196F3"  # 蓝色
        else:
            color = "#4CAF50"  # 绿色
            
        # 插入带颜色的文本
        self.chat_area.insert(tk.END, f"\n[{timestamp}] ", "timestamp")
        self.chat_area.insert(tk.END, f"{sender}：\n", "sender")
        self.chat_area.insert(tk.END, f"{message}\n", "message")
        
        # 设置标签颜色
        self.chat_area.tag_configure("timestamp", foreground="#999999")
        self.chat_area.tag_configure("sender", foreground=color, font=('Microsoft YaHei', 10, 'bold'))
        self.chat_area.tag_configure("message", foreground="#333333")
        
        self.chat_area.see(tk.END)
        self.chat_area.config(state='disabled')

    def _toggle_tts(self):
        """切换语音播报开关"""
        is_enabled = self.bot.toggle_tts()
        self.tts_button.config(text=f"语音播报: {'开' if is_enabled else '关'}")

    def _disable_input(self, disabled):
        """禁用/启用输入组件"""
        state = 'disabled' if disabled else 'normal'
        self.input_text.config(state=state)
        self.send_button.config(state=state)
        self.tts_button.config(state=state)
        self.camera_button.config(state=state)

    def _on_text_input(self, event=None):
        """处理文本输入"""
        # 如果正在语音播报，先停止
        if self.bot.is_speaking:
            self.bot.stop_processing()
            
        text = self.input_text.get().strip()
        
        # 如果有待处理的文件，将文件内容和提示词一起发送
        if self.pending_files:
            if not text:
                text = "请分析这些文件的内容"
            
            # 构建完整的提示词
            files_info = "\n".join([
                f"文件 {i+1}: {file['name']} ({file['type']})\n{file['content']}"
                for i, file in enumerate(self.pending_files)
            ])
            full_prompt = f"{text}\n\n文件内容：\n{files_info}"
            
            # 清空待处理文件列表
            self.pending_files = []
            self._update_file_status("")
            
            # 清空输入框并显示消息
            self.queue.put(lambda: self.input_text.delete(0, tk.END))
            self._thread_safe_display("您", text)
            
            # 启动处理线程
            threading.Thread(
                target=self._process_query,
                args=(full_prompt,),
                daemon=True
            ).start()
            return
        
        # 如果没有待处理文件，按原来的方式处理
        if not text:
            return
        
        # 清空输入框并显示消息
        self.queue.put(lambda: self.input_text.delete(0, tk.END))
        self._thread_safe_display("您", text)
        
        # 启动处理线程
        threading.Thread(
            target=self._process_query,
            args=(text,),
            daemon=True
        ).start()

    def _process_query(self, query):
        """处理用户查询"""
        try:
            self.queue.put(lambda: self._disable_input(True))
            # 使用text_chat方法处理响应
            response = self.bot.text_chat(query)
            self._thread_safe_display("AI", response)
            # 语音播报会在text_chat方法中自动处理
        except Exception as e:
            self._thread_safe_display("系统", f"处理错误: {str(e)}")
        finally:
            self.queue.put(lambda: self._disable_input(False))

    def _toggle_camera(self):
        """切换摄像头状态"""
        if self.bot.toggle_camera():
            self.camera_button.configure(text="关闭摄像头")
            self._display_message("系统", "摄像头已开启")
            # 启动摄像头画面更新
            self._update_camera_display()
        else:
            self.camera_button.configure(text="开启摄像头")
            self._display_message("系统", "摄像头已关闭")
            # 停止摄像头画面更新
            self.after_cancel(self._camera_update_id)
            # 清除摄像头显示
            self.camera_display.configure(image='')
            self.camera_display.image = None  # 清除引用

    def _update_camera_display(self):
        """更新摄像头画面显示"""
        if self.bot.is_camera_active and self.bot.camera:
            try:
                ret, frame = self.bot.camera.read()
                if ret:
                    # 调整图像大小以适应显示区域
                    height, width = frame.shape[:2]
                    max_size = 250  # 最大显示尺寸
                    scale = min(max_size/width, max_size/height)
                    new_width = int(width * scale)
                    new_height = int(height * scale)
                    frame = cv2.resize(frame, (new_width, new_height))
                    
                    # 转换颜色空间并显示
                    frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                    image = Image.fromarray(frame_rgb)
                    photo = ImageTk.PhotoImage(image=image)
                    self.camera_display.configure(image=photo)
                    self.camera_display.image = photo  # 保持引用
            except Exception as e:
                print(f"摄像头画面更新错误: {str(e)}")
        
        # 继续更新
        self._camera_update_id = self.after(30, self._update_camera_display)  # 约30fps

    def _open_file_dialog(self):
        """打开文件选择对话框"""
        filetypes = (
            ('所有支持的文件', '*.pdf;*.docx;*.pptx;*.png;*.jpg;*.jpeg'),
            ('PDF文件', '*.pdf'),
            ('Word文件', '*.docx'),
            ('PPT文件', '*.pptx'),
            ('图片文件', '*.png;*.jpg;*.jpeg'),
            ('所有文件', '*.*')
        )
        
        filename = filedialog.askopenfilename(
            title='选择文件',
            filetypes=filetypes
        )
        
        if filename:
            self._process_file(filename)

    def _handle_drop(self, event):
        """处理文件拖放"""
        file_path = event.data
        # 移除可能的大括号
        file_path = file_path.strip('{}')
        self._process_file(file_path)

    def _process_file(self, file_path):
        """处理上传的文件"""
        try:
            file_name = os.path.basename(file_path)
            file_ext = os.path.splitext(file_path)[1].lower()
            
            # 更新状态为处理中
            self._update_file_status(f"正在处理文件: {file_name}")
            
            # 创建处理线程
            def process_thread():
                try:
                    if file_ext in ['.pdf']:
                        content = self._extract_pdf_content(file_path)
                    elif file_ext in ['.docx']:
                        content = self._extract_docx_content(file_path)
                    elif file_ext in ['.pptx']:
                        content = self._extract_pptx_content(file_path)
                    elif file_ext in ['.png', '.jpg', '.jpeg']:
                        content = self._process_image(file_path)
                    else:
                        raise ValueError(f"不支持的文件类型: {file_ext}")
                    
                    # 将处理结果添加到待处理列表
                    self.pending_files.append({
                        'name': file_name,
                        'content': content,
                        'type': file_ext
                    })
                    
                    # 更新状态为处理完成
                    self.queue.put(lambda: self._update_file_status(
                        f"文件处理完成: {file_name}\n"
                        f"请在输入框输入提示词（可选）后点击发送"
                    ))
                    
                except Exception as e:
                    self.queue.put(lambda: self._update_file_status(
                        f"文件处理错误: {file_name}\n{str(e)}"
                    ))
            
            # 启动处理线程
            threading.Thread(target=process_thread, daemon=True).start()
            
        except Exception as e:
            self._update_file_status(f"文件处理错误: {str(e)}")

    def _extract_pdf_content(self, file_path):
        """提取PDF文件内容"""
        doc = fitz.open(file_path)
        text = ""
        for page in doc:
            text += page.get_text()
        doc.close()
        return text

    def _extract_docx_content(self, file_path):
        """提取Word文件内容"""
        doc = Document(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text

    def _extract_pptx_content(self, file_path):
        """提取PPT文件内容"""
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    def _process_image(self, file_path):
        """处理图片文件"""
        try:
            # 读取图片并转换为base64
            with Image.open(file_path) as img:
                buffered = BytesIO()
                img.save(buffered, format="PNG")
                img_str = base64.b64encode(buffered.getvalue()).decode()
            
            # 调用豆包视觉API分析图片
            return self.bot._call_doubao_vision_api(img_str, "请分析这张图片的内容")
            
        except Exception as e:
            raise Exception(f"图片处理错误: {str(e)}")

    def _update_file_status(self, message):
        """更新文件处理状态"""
        self.file_status_label.configure(text=message)